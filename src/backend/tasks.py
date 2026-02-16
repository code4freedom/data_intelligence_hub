import os
import json
import re
from pathlib import Path
from typing import Optional
from datetime import datetime
import uuid
from jinja2 import Template
from pptx import Presentation
from pptx.util import Inches
from src.backend.intelligence import compute_manifest_intelligence
from src.backend.advanced_analytics import compute_advanced_analytics

# pyppeteer for headless rendering
import asyncio
from pyppeteer import launch

# reportlab for error fallback
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib.pagesizes import letter

DATA_DIR = Path('/data')
EXPORTS_DIR = DATA_DIR / 'exports'
EXPORTS_DIR.mkdir(parents=True, exist_ok=True)
TEMPLATES_DIR = Path('/app/frontend/templates')
MANIFESTS_DIR = DATA_DIR / 'manifests'
CHUNKS_DIR = DATA_DIR / 'chunks'
APPENDIX_DIR = TEMPLATES_DIR / 'appendices'
APPENDIX_CATALOG = [
    {"id": "current_state_architecture", "title": "Current State Architecture", "file": "current_state_architecture_template.html"},
    {"id": "future_state_architecture", "title": "Future State Architecture", "file": "future_state_architecture_template.html"},
    {"id": "application_landscape", "title": "Application Landscape", "file": "application_landscape_template.html"},
    {"id": "vsphere_estate_platform_health", "title": "vSphere Estate - Platform Health", "file": "platform_health_template.html"},
]


def _appendix_fragment_from_html(html_doc: str) -> str:
    """Convert a full HTML doc to an embeddable fragment with local styles."""
    if not html_doc:
        return ""
    style_blocks = re.findall(r"<style[^>]*>.*?</style>", html_doc, flags=re.IGNORECASE | re.DOTALL)
    body_match = re.search(r"<body[^>]*>(.*)</body>", html_doc, flags=re.IGNORECASE | re.DOTALL)
    body_content = body_match.group(1) if body_match else html_doc
    return "\n".join(style_blocks) + "\n" + body_content


def _resolve_context_dirs(manifest_path: Path):
    """Resolve manifest/chunk/export dirs for project-based and legacy layouts."""
    # Project layout: /data/projects/<project>/manifests/<manifest>
    if "projects" in manifest_path.parts and "manifests" in manifest_path.parts:
        manifests_dir = manifest_path.parent
        project_base = manifests_dir.parent
        chunks_dir = project_base / "chunks"
        exports_dir = project_base / "exports"
        exports_dir.mkdir(parents=True, exist_ok=True)
        return manifests_dir, chunks_dir, exports_dir
    # Legacy fallback.
    return MANIFESTS_DIR, CHUNKS_DIR, EXPORTS_DIR


def _resolve_project_app_mapping(manifest_path: Path) -> Optional[str]:
    if "projects" in manifest_path.parts and "manifests" in manifest_path.parts:
        manifests_dir = manifest_path.parent
        project_base = manifests_dir.parent
        p = project_base / "config" / "app_mapping.csv"
        if p.exists():
            return str(p)
    env_default = os.environ.get("APP_MAPPING_CSV")
    return env_default if env_default else None


async def _render_html_assets(html: str, out_pdf: Optional[Path] = None, out_png: Optional[Path] = None):
    browser = None
    try:
        browser = await launch(options={
            'args': ['--no-sandbox', '--disable-setuid-sandbox', '--disable-dev-shm-usage'],
            'headless': True,
            'handleSIGINT': False,
            'handleSIGTERM': False,
            'handleSIGHUP': False,
        })
        page = await browser.newPage()
        
        # Use a wide viewport to match landscape-oriented report pages.
        await page.setViewport({'width': 1600, 'height': 980})
        
        # Load HTML (no waitUntil parameter - not available for setContent)
        await page.setContent(html)
        
        # Wait a bit for styles to apply
        await asyncio.sleep(0.8)

        if out_pdf:
            pdf_options = {
                'path': str(out_pdf),
                'format': 'A4',
                'landscape': True,
                'printBackground': True,
                'margin': {
                    'top': '10mm',
                    'bottom': '10mm',
                    'left': '10mm',
                    'right': '10mm'
                }
            }
            await page.pdf(pdf_options)

        if out_png:
            await page.screenshot({'path': str(out_png), 'fullPage': True})

    except Exception as e:
        print(f"Error in pyppeteer rendering: {e}")
        if out_pdf:
            c = rl_canvas.Canvas(str(out_pdf), pagesize=letter)
            c.setFont("Helvetica-Bold", 14)
            c.drawString(50, 750, "Report Generation Error")
            c.setFont("Helvetica", 10)
            c.drawString(50, 700, f"Error: {str(e)[:100]}")
            c.save()
        raise
    finally:
        if browser:
            await browser.close()


def generate_pptx_from_png(png_path: Path, out_pptx: Path):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = Inches(0.5)
    pic = slide.shapes.add_picture(str(png_path), left, top, width=prs.slide_width - Inches(1))
    prs.save(str(out_pptx))


def generate_report(
    manifest_path: str,
    template_name: str = 'vsphere',
    output_format: str = 'pdf',
    appendices: Optional[object] = None,
) -> dict:
    """Task: read manifest, render template, produce report outputs, return export paths."""
    manifest_p = Path(manifest_path)
    if not manifest_p.exists():
        raise FileNotFoundError(manifest_path)
    manifests_dir, chunks_dir, exports_dir = _resolve_context_dirs(manifest_p)
    manifest = json.loads(manifest_p.read_text(encoding='utf-8'))

    # Compute KPIs ONLY from chunks referenced in this manifest, not all chunks
    ingest_id = manifest.get('ingest_id', 'ingest')
    sheet = manifest.get('sheet', 'vInfo')
    chunk_count = manifest.get('chunk_count', 0)
    total_vms = manifest.get('total_rows', 0)
    
    # Compute metrics from this manifest's chunks only
    total_hosts = set()
    total_compute = 0
    total_memory_mb = 0
    eos_risk = 0
    from src.schema.rvtools_schema import find_column
    
    for chunk_info in manifest.get('chunks', []):
        chunk_path = Path(chunk_info.get('local_path'))
        if not chunk_path.exists():
            continue
        
        try:
            import pandas as pd
            df = pd.read_parquet(str(chunk_path))
            
            # Host count
            host_col = find_column(df.columns, 'vInfo', 'host')
            if host_col and host_col in df.columns:
                host_vals = df[host_col].dropna().unique().tolist()
                total_hosts.update([str(h) for h in host_vals])
            
            # CPU
            cpu_col = find_column(df.columns, 'vInfo', 'numcpu')
            if cpu_col and cpu_col in df.columns:
                try:
                    total_compute += int(df[cpu_col].fillna(0).astype(int).sum())
                except Exception:
                    pass
            
            # Memory
            mem_col = find_column(df.columns, 'vInfo', 'memorymb')
            if mem_col and mem_col in df.columns:
                try:
                    total_memory_mb += df[mem_col].fillna(0).astype(float).sum()
                except Exception:
                    pass
            
            # EOS Risk (versions starting with 6 or 7 are older)
            ver_col = find_column(df.columns, 'vInfo', 'version')
            if ver_col and ver_col in df.columns:
                try:
                    versions = df[ver_col].astype(str).str[0]
                    eos_risk += (versions.isin(['6', '7'])).sum()
                except Exception:
                    pass
        except Exception as e:
            print(f"Warning: could not read chunk {chunk_path}: {e}")
    
    total_memory_tb = total_memory_mb / (1024 * 1024) if total_memory_mb > 0 else 0
    
    context = {
        'total_vms': total_vms,
        'total_hosts': len(total_hosts),
        'total_compute': int(total_compute),
        'total_memory': f"{total_memory_tb:.2f} TB" if total_memory_tb > 0 else "--",
        'eos_risk': int(eos_risk),
        'chunks': chunk_count,
        'generation_date': datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
        'author_name': 'Samir Roshan',
        'logo_url': f"{os.environ.get('REPORT_ASSET_BASE_URL', 'http://localhost:8000')}/static/public/vcf-hub-logo.png",
    }

    selected_ids = []
    if appendices is None:
        selected_ids = [a["id"] for a in APPENDIX_CATALOG]
    elif isinstance(appendices, str):
        selected_ids = [x.strip() for x in appendices.split(",") if x.strip()]
    elif isinstance(appendices, list):
        selected_ids = [str(x).strip() for x in appendices if str(x).strip()]

    allowed = {a["id"]: a for a in APPENDIX_CATALOG}
    example_pages = []
    for aid in selected_ids:
        meta = allowed.get(aid)
        if not meta:
            continue
        p = APPENDIX_DIR / meta["file"]
        if not p.exists():
            continue
        try:
            tpl_text = p.read_text(encoding='utf-8')
            html_doc = Template(tpl_text).render(**context)
            html_fragment = _appendix_fragment_from_html(html_doc)
            example_pages.append({"id": aid, "title": meta["title"], "html": html_fragment})
        except Exception:
            continue
    context["example_pages"] = example_pages

    intelligence = compute_manifest_intelligence(
        manifest_path,
        str(manifests_dir),
        str(chunks_dir),
        app_map_csv=_resolve_project_app_mapping(manifest_p),
    )
    context.update(
        {
            "executive_score": intelligence.get("executive_score", 0),
            "executive_components": intelligence.get("executive_components", {}),
            "consolidation": intelligence.get("consolidation", {}),
            "lifecycle": intelligence.get("lifecycle", {}),
            "performance": intelligence.get("performance", {}),
            "application": intelligence.get("application", {}),
            "logical_view": intelligence.get("logical_view", {}),
            "trend_growth": intelligence.get("trend_growth", {}),
            "trends": intelligence.get("trends", []),
            "insights": intelligence.get("insights", []),
            "top_clusters": intelligence.get("consolidation", {}).get("cluster_density_top", []),
        }
    )
    advanced = compute_advanced_analytics(
        manifest_path,
        str(manifests_dir),
        str(chunks_dir),
    )
    context.update(
        {
            "forecasting": advanced.get("forecasting", {}),
            "anomalies": advanced.get("anomalies", {}),
            "right_sizing_recommendations": advanced.get("right_sizing_recommendations", []),
            "eos_prioritization": advanced.get("eos_prioritization", []),
            "consolidation_optimization": advanced.get("consolidation_optimization", {}),
            "dependency_graph": advanced.get("dependency_graph", {}),
            "storage_efficiency": advanced.get("storage_efficiency", {}),
            "drift_governance": advanced.get("drift_governance", {}),
            "what_if_simulation": advanced.get("what_if_simulation", {}),
            "operational_scorecard": advanced.get("operational_scorecard", {}),
            "action_backlog_ml": advanced.get("action_backlog", []),
        }
    )

    # load template
    tpl_path = TEMPLATES_DIR / f"{template_name}_template.html"
    if not tpl_path.exists():
        raise FileNotFoundError(str(tpl_path))
    tpl_text = tpl_path.read_text(encoding='utf-8')
    rendered = Template(tpl_text).render(**context)

    # output files
    ingest_id = manifest.get('ingest_id', 'ingest')
    out_base = exports_dir / ingest_id
    out_base.mkdir(parents=True, exist_ok=True)
    run_id = f"{datetime.utcnow().strftime('%Y%m%d%H%M%S')}_{uuid.uuid4().hex[:8]}"
    out_pdf = out_base / f"report_{template_name}_{ingest_id}_{run_id}.pdf"
    out_png = out_base / f"report_{template_name}_{ingest_id}_{run_id}.png"
    out_pptx = out_base / f"report_{template_name}_{ingest_id}_{run_id}.pptx"
    fmt = (output_format or 'pdf').lower()
    include_pdf = fmt in ('pdf', 'both')
    include_pptx = fmt in ('pptx', 'both')

    if not include_pdf and not include_pptx:
        raise ValueError(f"Unsupported output_format: {output_format}")

    needs_png = include_pptx
    asyncio.run(
        _render_html_assets(
            rendered,
            out_pdf=out_pdf if include_pdf else None,
            out_png=out_png if needs_png else None,
        )
    )

    if include_pptx:
        generate_pptx_from_png(out_png, out_pptx)

    result = {
        'ingest_id': ingest_id,
        'chunks': chunk_count,
        'output_format': fmt,
        'intelligence': intelligence,
        'advanced': advanced,
    }
    if include_pdf:
        result['pdf'] = str(out_pdf)
    if include_pptx:
        result['pptx'] = str(out_pptx)
    if needs_png:
        result['png'] = str(out_png)
    return result


if __name__ == '__main__':
    # quick test harness: pass manifest path via env
    mp = os.environ.get('MANIFEST_PATH')
    if not mp:
        print('Set MANIFEST_PATH for local test')
    else:
        print(generate_report(mp))
